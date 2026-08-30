"""Office documents for the extraction tests, written by real authoring libraries.

A container assembled by hand proves less than it looks: it encodes the same
assumptions the extractor makes, so the two agree with each other rather than
with Word. python-docx, openpyxl, python-pptx and odfpy emit the part layout,
namespaces and scaffolding a word processor actually produces, which is what
the extractor has to be tested against.

They are dev-only and only ever write here; nothing under
workspace/common/documents imports them. The degenerate cases have no
authoring library, because no library produces a part that is not XML, a
member that lies about its size, or an entity bomb - make_zip() stays for
those, and for the shapes a library refuses to emit.
"""

from __future__ import annotations

import io
import zipfile

from docx import Document
from odf.draw import Frame, Page, TextBox
from odf.opendocument import (
    OpenDocumentPresentation,
    OpenDocumentSpreadsheet,
    OpenDocumentText,
)
from odf.table import Table, TableCell, TableRow
from odf.text import H, P
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

ODT = "application/vnd.oasis.opendocument.text"
ODS = "application/vnd.oasis.opendocument.spreadsheet"
ODP = "application/vnd.oasis.opendocument.presentation"

# The layout with no placeholders on it, so a slide holds the text the test put
# there and nothing the template came with.
_BLANK_SLIDE_LAYOUT = 6


def make_zip(members: dict[str, str | bytes], *, first: str | None = None) -> bytes:
    """Zip *members* verbatim. *first* is stored uncompressed, as ODF wants.

    For the containers no authoring library will produce.
    """
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        if first is not None:
            archive.writestr(zipfile.ZipInfo(first), members[first])
        for name, payload in members.items():
            if name != first:
                archive.writestr(name, payload)
    return buffer.getvalue()


def _with_extra(payload: bytes, extra: dict[str, str] | None) -> bytes:
    """Append parts an authoring library cannot write, such as footnotes."""
    if not extra:
        return payload
    buffer = io.BytesIO(payload)
    with zipfile.ZipFile(buffer, "a", zipfile.ZIP_DEFLATED) as archive:
        for name, content in extra.items():
            archive.writestr(name, content)
    return buffer.getvalue()


def make_docx(
    paragraphs: list[str],
    *,
    heading: str | None = None,
    runs: list[str] | None = None,
    table: list[list[str]] | None = None,
    header: str | None = None,
    footer: str | None = None,
    extra: dict[str, str] | None = None,
) -> bytes:
    """A .docx as Word lays one out. *runs* build one paragraph from many runs."""
    document = Document()
    if heading is not None:
        document.add_heading(heading, level=1)
    for text in paragraphs:
        document.add_paragraph(text)
    if runs:
        paragraph = document.add_paragraph()
        for run in runs:
            paragraph.add_run(run)
    if table:
        added = document.add_table(rows=len(table), cols=len(table[0]))
        for row, cells in zip(added.rows, table, strict=True):
            for cell, text in zip(row.cells, cells, strict=True):
                cell.text = text
    section = document.sections[0]
    if header is not None:
        section.header.paragraphs[0].text = header
    if footer is not None:
        section.footer.paragraphs[0].text = footer

    buffer = io.BytesIO()
    document.save(buffer)
    return _with_extra(buffer.getvalue(), extra)


def make_xlsx(*, sheets: dict[str, list[list]] | None = None) -> bytes:
    """A workbook, one entry per sheet, each a list of rows.

    openpyxl writes text inline, as <c t="inlineStr"><is><t>, and never emits
    a sharedStrings part. That is one of the two dialects the extractor has to
    read; make_xlsx_shared_strings() below is the other.
    """
    workbook = Workbook()
    workbook.remove(workbook.active)
    for title, rows in (sheets or {"Sheet1": []}).items():
        worksheet = workbook.create_sheet(title)
        for row in rows:
            worksheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def make_pptx(slides: list[list[str]]) -> bytes:
    """A deck, one text box per slide, one paragraph per string."""
    presentation = Presentation()
    layout = presentation.slide_layouts[_BLANK_SLIDE_LAYOUT]
    for paragraphs in slides:
        slide = presentation.slides.add_slide(layout)
        frame = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(2)
        ).text_frame
        frame.text = paragraphs[0]
        for text in paragraphs[1:]:
            frame.add_paragraph().text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_odf(
    mime_type: str, paragraphs: list[str], *, heading: str | None = None
) -> bytes:
    """An OpenDocument container; odt, ods and odp differ in where text sits."""
    buffer = io.BytesIO()
    if mime_type == ODT:
        document = OpenDocumentText()
        if heading is not None:
            document.text.addElement(H(outlinelevel=1, text=heading))
        for text in paragraphs:
            document.text.addElement(P(text=text))
    elif mime_type == ODS:
        document = OpenDocumentSpreadsheet()
        table = Table(name="Sheet1")
        for text in paragraphs:
            row, cell = TableRow(), TableCell()
            cell.addElement(P(text=text))
            row.addElement(cell)
            table.addElement(row)
        document.spreadsheet.addElement(table)
    elif mime_type == ODP:
        document = OpenDocumentPresentation()
        page = Page(name="page1", masterpagename="Default")
        frame = Frame(width="20cm", height="4cm", x="1cm", y="1cm")
        box = TextBox()
        for text in paragraphs:
            box.addElement(P(text=text))
        frame.addElement(box)
        page.addElement(frame)
        document.presentation.addElement(page)
    else:
        raise ValueError(f"Not an OpenDocument type: {mime_type}")
    document.save(buffer)
    return buffer.getvalue()


# Excel keeps a cell's text in xl/sharedStrings.xml and leaves only a row
# number in the cell; openpyxl writes the text inline instead. Both are valid
# and both turn up in a real workspace, so both paths need a fixture - and
# openpyxl will not emit this one.
_SPREADSHEETML = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"


def make_xlsx_shared_strings(strings: list[str], *, numbers: list[str] = ()) -> bytes:
    """A workbook in Excel's dialect: cells of type "s" indexing sharedStrings."""
    items = "".join(f"<si><t>{text}</t></si>" for text in strings)
    cells = [f'<c t="s"><v>{index}</v></c>' for index in range(len(strings))]
    cells += [f"<c><v>{value}</v></c>" for value in numbers]
    return make_zip(
        {
            "xl/sharedStrings.xml": (
                f'<sst xmlns="{_SPREADSHEETML}" count="{len(strings)}" '
                f'uniqueCount="{len(strings)}">{items}</sst>'
            ),
            "xl/worksheets/sheet1.xml": (
                f'<worksheet xmlns="{_SPREADSHEETML}"><sheetData>'
                f'<row r="1">{"".join(cells)}</row></sheetData></worksheet>'
            ),
        }
    )
